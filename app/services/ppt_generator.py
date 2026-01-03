"""Production-grade PowerPoint generation service using python-pptx.

Designed to match McKinsey, BCG, Bain, and JPMorgan visual standards.
Incorporates master templates, professional design, and investment banking aesthetics.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from typing import List, Dict, Any, Optional
import plotly.graph_objects as go
from io import BytesIO


class PPTGenerator:
    """Generate McKinsey/BCG/JPM-grade PowerPoint presentations."""
    
    # Professional color palettes (matching PDF Generator)
    BRAND_COLORS = {
        'mckinsey': {
            'primary': RGBColor(0, 63, 92),
            'accent': RGBColor(47, 75, 124),
            'highlight': RGBColor(249, 93, 106),
            'text_dark': RGBColor(26, 26, 26),
            'text_light': RGBColor(102, 102, 102)
        },
        'bcg': {
            'primary': RGBColor(0, 167, 88),
            'accent': RGBColor(0, 63, 92),
            'highlight': RGBColor(255, 111, 60),
            'text_dark': RGBColor(0, 0, 0),
            'text_light': RGBColor(74, 74, 74)
        },
        'jpmorgan': {
            'primary': RGBColor(0, 94, 184),
            'accent': RGBColor(51, 51, 51),
            'highlight': RGBColor(212, 175, 55),
            'text_dark': RGBColor(0, 0, 0),
            'text_light': RGBColor(102, 102, 102)
        },
        'bain': {
            'primary': RGBColor(237, 28, 36),
            'accent': RGBColor(31, 41, 55),
            'highlight': RGBColor(245, 158, 11),
            'text_dark': RGBColor(17, 24, 39),
            'text_light': RGBColor(107, 114, 128)
        }
    }
    
    def __init__(self, brand: str = 'mckinsey'):
        """Initialize PPT generator with brand-specific styling."""
        self.brand = brand.lower()
        self.colors = self.BRAND_COLORS.get(self.brand, self.BRAND_COLORS['mckinsey'])
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    async def generate_ppt(
        self,
        slides: List[Dict[str, Any]],
        output_path: str,
        title: Optional[str] = None,
        subtitle: Optional[str] = None,
        company_name: Optional[str] = None
    ) -> str:
        """Generate PowerPoint presentation with professional design."""
        try:
            # Add cover slide if title provided
            if title:
                self._add_cover_slide(title, subtitle, company_name)
            
            # Generate each slide
            for slide_data in slides:
                slide_type = slide_data.get('type', 'content')
                
                if slide_type == 'title':
                    self._add_title_slide(slide_data)
                elif slide_type == 'section_divider':
                    self._add_section_divider_slide(slide_data)
                elif slide_type == 'chart':
                    self._add_chart_slide(slide_data)
                elif slide_type == 'two_column':
                    self._add_two_column_slide(slide_data)
                else:
                    self._add_content_slide(slide_data)
            
            # Add closing slide
            if company_name:
                self._add_closing_slide(company_name)
            
            self.prs.save(output_path)
            return output_path
            
        except Exception as e:
            print(f"PPT generation failed: {e}")
            raise
    
    def _add_cover_slide(self, title: str, subtitle: Optional[str], company_name: Optional[str]):
        """Add professional cover slide with gradient background."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Brand accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(2.5), Inches(8), Inches(0.15)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.colors['primary']
        accent_bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            
            sub_para = subtitle_frame.paragraphs[0]
            sub_para.font.size = Pt(28)
            sub_para.font.color.rgb = self.colors['text_light']
            sub_para.alignment = PP_ALIGN.CENTER
        
        # Company name
        if company_name:
            company_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
            company_frame = company_box.text_frame
            company_frame.text = company_name
            
            company_para = company_frame.paragraphs[0]
            company_para.font.size = Pt(20)
            company_para.font.bold = True
            company_para.font.color.rgb = self.colors['text_dark']
            company_para.alignment = PP_ALIGN.CENTER
    
    def _add_section_divider_slide(self, slide_data: Dict):
        """Add section divider slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Colored background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['primary']
        
        # Section title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', '')
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
    
    def _add_title_slide(self, slide_data: Dict[str, Any]):
        """Add title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header_line(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', '')
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if slide_data.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = slide_data['subtitle']
            
            sub_para = subtitle_frame.paragraphs[0]
            sub_para.font.size = Pt(24)
            sub_para.font.color.rgb = self.colors['text_light']
            sub_para.alignment = PP_ALIGN.CENTER
        
        self._add_footer(slide)
    
    def _add_content_slide(self, slide_data: Dict[str, Any]):
        """Add content slide with smart bullet formatting."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header_line(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', '')
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']
        
        # Content with smart bullets
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        for idx, item in enumerate(slide_data.get('content', [])):
            if isinstance(item, str) and item.strip():
                level = 0
                clean_item = item.strip()
                
                if clean_item.startswith('  -') or clean_item.startswith('    •'):
                    level = 1
                    clean_item = clean_item.lstrip(' -•').strip()
                else:
                    clean_item = clean_item.lstrip('•-').strip()
                
                clean_item = clean_item.replace('**', '')
                
                p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
                p.text = clean_item
                p.level = level
                p.font.size = Pt(18) if level == 0 else Pt(16)
                p.font.color.rgb = self.colors['text_dark']
                p.space_after = Pt(14) if level == 0 else Pt(10)
        
        if slide_data.get('speaker_notes'):
            self._add_speaker_notes(slide, slide_data['speaker_notes'])
        
        self._add_footer(slide)
    
    def _add_two_column_slide(self, slide_data: Dict):
        """Add two-column layout slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header_line(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', '')
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        left_frame.clear()
        
        for idx, item in enumerate(slide_data.get('left_column', [])):
            if isinstance(item, str):
                p = left_frame.add_paragraph() if idx > 0 else left_frame.paragraphs[0]
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = self.colors['text_dark']
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.5), Inches(5.5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        right_frame.clear()
        
        for idx, item in enumerate(slide_data.get('right_column', [])):
            if isinstance(item, str):
                p = right_frame.add_paragraph() if idx > 0 else right_frame.paragraphs[0]
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = self.colors['text_dark']
        
        self._add_footer(slide)
    
    def _add_chart_slide(self, slide_data: Dict[str, Any]):
        """Add slide with high-quality chart."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header_line(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', '')
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']
        
        # Chart
        chart_data = slide_data.get('chart_data')
        if chart_data:
            chart_img = self._plotly_to_bytes(chart_data)
            if chart_img:
                slide.shapes.add_picture(
                    chart_img,
                    Inches(0.5),
                    Inches(1.3),
                    width=Inches(6.5),
                    height=Inches(4.5)
                )
        
        # Key insights
        insights_box = slide.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(2.5), Inches(4.5))
        text_frame = insights_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        for idx, item in enumerate(slide_data.get('content', [])):
            if isinstance(item, str) and item.strip():
                clean_item = item.replace('•', '').strip()
                p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
                p.text = f"• {clean_item}"
                p.font.size = Pt(12)
                p.font.color.rgb = self.colors['text_dark']
        
        if slide_data.get('speaker_notes'):
            self._add_speaker_notes(slide, slide_data['speaker_notes'])
        
        self._add_footer(slide)
    
    def _add_closing_slide(self, company_name: str):
        """Add professional closing slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Thank you message
        thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        thank_you_frame = thank_you_box.text_frame
        thank_you_frame.text = "Thank You"
        
        thank_you_para = thank_you_frame.paragraphs[0]
        thank_you_para.font.size = Pt(54)
        thank_you_para.font.bold = True
        thank_you_para.font.color.rgb = self.colors['primary']
        thank_you_para.alignment = PP_ALIGN.CENTER
        
        # Company name
        company_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
        company_frame = company_box.text_frame
        company_frame.text = company_name
        
        company_para = company_frame.paragraphs[0]
        company_para.font.size = Pt(24)
        company_para.font.color.rgb = self.colors['text_light']
        company_para.alignment = PP_ALIGN.CENTER
    
    def _add_header_line(self, slide):
        """Add brand-colored header accent line."""
        header_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.05)
        )
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = self.colors['primary']
        header_line.line.fill.background()
    
    def _add_footer(self, slide):
        """Add footer with confidential mark."""
        footer_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.45), Inches(10), Inches(0.05)
        )
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = self.colors['text_light']
        footer_line.line.fill.background()
        
        conf_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(3), Inches(0.3))
        conf_frame = conf_box.text_frame
        conf_frame.text = "CONFIDENTIAL"
        
        conf_para = conf_frame.paragraphs[0]
        conf_para.font.size = Pt(9)
        conf_para.font.color.rgb = self.colors['text_light']
    
    def _add_speaker_notes(self, slide, notes: str):
        """Add enhanced speaker notes."""
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes
    
    def _plotly_to_bytes(self, chart_data: Dict) -> BytesIO:
        """Convert Plotly chart to high-quality image bytes."""
        try:
            fig = go.Figure(chart_data)
            
            fig.update_layout(
                plot_bgcolor='white',
                paper_bgcolor='white',
                font=dict(family='Helvetica', size=12),
                margin=dict(l=60, r=60, t=60, b=60)
            )
            
            img_bytes = fig.to_image(format="png", width=1300, height=900, scale=2)
            return BytesIO(img_bytes)
        except Exception as e:
            print(f"Chart conversion failed: {e}")
            return None
