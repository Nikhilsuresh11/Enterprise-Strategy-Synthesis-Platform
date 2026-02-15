"""Pitch Deck Generation Service - Creates PowerPoint presentations from analysis."""

from typing import Dict, Any, List
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from app.utils.logger import get_logger

logger = get_logger(__name__)


class PitchDeckService:
    """
    Service for generating professional pitch decks from strategic analysis.
    
    Features:
    - 10-slide professional template
    - 3-4 bullet points per slide
    - Clean, modern design
    - Automatic formatting
    """
    
    def __init__(self):
        self.title_font_size = Pt(44)
        self.heading_font_size = Pt(32)
        self.body_font_size = Pt(18)
        
        # Brand colors (professional blue theme)
        self.primary_color = RGBColor(0, 51, 102)      # Dark blue
        self.secondary_color = RGBColor(0, 102, 204)   # Medium blue
        self.accent_color = RGBColor(255, 153, 0)      # Orange
        
        logger.info("pitch_deck_service_initialized")
    
    async def generate_deck(
        self,
        company_name: str,
        analysis_data: Dict[str, Any],
        output_path: str
    ) -> str:
        """
        Generate pitch deck from analysis data.
        
        Args:
            company_name: Name of the company
            analysis_data: Complete analysis from all agents
            output_path: Path to save the .pptx file
        
        Returns:
            Path to generated deck
        """
        logger.info("generating_pitch_deck", company=company_name)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Generate slides
        self._add_title_slide(prs, company_name, analysis_data)
        self._add_company_overview_slide(prs, company_name, analysis_data)
        self._add_market_opportunity_slide(prs, analysis_data)
        self._add_competitive_landscape_slide(prs, analysis_data)
        self._add_financial_highlights_slide(prs, analysis_data)
        self._add_risk_assessment_slide(prs, analysis_data)
        self._add_strategic_recommendations_slide(prs, analysis_data)
        self._add_swot_slide(prs, analysis_data)
        self._add_next_steps_slide(prs, analysis_data)
        self._add_qa_slide(prs)
        
        # Save presentation
        prs.save(output_path)
        
        logger.info("pitch_deck_generated", path=output_path, slides=len(prs.slides))
        
        return output_path
    
    def _add_title_slide(self, prs: Presentation, company_name: str, data: Dict[str, Any]):
        """Slide 1: Title slide with company name and tagline."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Company name
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"{company_name}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = self.title_font_size
        title_para.font.bold = True
        title_para.font.color.rgb = self.primary_color
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.8), Inches(8), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Strategic Analysis & Recommendations"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = self.secondary_color
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(8), Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime("%B %Y")
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = RGBColor(128, 128, 128)
        date_para.alignment = PP_ALIGN.CENTER
    
    def _add_company_overview_slide(self, prs: Presentation, company_name: str, data: Dict[str, Any]):
        """Slide 2: Company Overview - 4-5 key facts."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        
        # Title
        title = slide.shapes.title
        title.text = "Company Overview"
        title.text_frame.paragraphs[0].font.size = self.heading_font_size
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        # Content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Get key facts from company profile
        company_profile = data.get("company_profile", {})
        key_facts = company_profile.get("key_facts", [])
        
        if not key_facts:
            key_facts = ["Company information unavailable"]
        
        # Add bullet points (max 4)
        for fact in key_facts[:4]:
            p = text_frame.add_paragraph()
            p.text = fact
            p.level = 0
            p.font.size = self.body_font_size
            p.space_after = Pt(12)
    
    def _add_market_opportunity_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Slide 3: Market Opportunity - 4-5 market insights."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Market Opportunity"
        title.text_frame.paragraphs[0].font.size = self.heading_font_size
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Get market insights
        market_analysis = data.get("market_analysis", {})
        insights = market_analysis.get("key_insights", [])
        
        if not insights:
            insights = ["Market analysis unavailable"]
        
        # Add bullet points (max 4)
        for insight in insights[:4]:
            p = text_frame.add_paragraph()
            p.text = insight
            p.level = 0
            p.font.size = self.body_font_size
            p.space_after = Pt(12)
    
    def _add_competitive_landscape_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Slide 4: Competitive Landscape - Top 3-4 competitors."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Competitive Landscape"
        title.text_frame.paragraphs[0].font.size = self.heading_font_size
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Get competitors
        competitors = data.get("competitor_analysis", [])
        
        if not competitors:
            p = text_frame.add_paragraph()
            p.text = "Competitor analysis unavailable"
            p.font.size = self.body_font_size
        else:
            # Add each competitor (max 4)
            for comp in competitors[:4]:
                name = comp.get("name", "Unknown")
                key_point = comp.get("key_point", "No details available")
                
                p = text_frame.add_paragraph()
                p.text = f"{name}: {key_point}"
                p.level = 0
                p.font.size = self.body_font_size
                p.space_after = Pt(12)
    
    def _add_financial_highlights_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Slide 5: Financial Highlights - 4-5 key metrics."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Financial Highlights"
        title.text_frame.paragraphs[0].font.size = self.heading_font_size
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Get financial highlights
        financial_model = data.get("financial_model", {})
        highlights = financial_model.get("key_highlights", [])
        
        if not highlights:
            highlights = ["Financial data unavailable"]
        
        # Add bullet points (max 4)
        for highlight in highlights[:4]:
            p = text_frame.add_paragraph()
            p.text = highlight
            p.level = 0
            p.font.size = self.body_font_size
            p.space_after = Pt(12)
    
    def _add_risk_assessment_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Slide 6: Risk Assessment - Top 4 risks."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Key Risks & Mitigation"
        title.text_frame.paragraphs[0].font.size = self.heading_font_size
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Get top risks
        risk_assessment = data.get("risk_assessment", {})
        risks = risk_assessment.get("top_risks", [])
        
        if not risks:
            p = text_frame.add_paragraph()
            p.text = "Risk assessment unavailable"
            p.font.size = self.body_font_size
        else:
            # Add each risk (max 4)
            for risk in risks[:4]:
                risk_name = risk.get("risk", "Unknown risk")
                severity = risk.get("severity", "unknown")
                description = risk.get("description", "")
                
                # Format: Risk name (severity): description
                p = text_frame.add_paragraph()
                p.text = f"{risk_name} ({severity.upper()}): {description}"
                p.level = 0
                p.font.size = self.body_font_size
                p.space_after = Pt(12)
    
    def _add_strategic_recommendations_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Slide 7: Strategic Recommendations - 4 key recommendations."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Strategic Recommendations"
        title.text_frame.paragraphs[0].font.size = self.heading_font_size
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Get recommendations
        strategy = data.get("strategy_synthesis", {})
        recommendations = strategy.get("key_recommendations", [])
        
        if not recommendations:
            recommendations = ["Strategic recommendations unavailable"]
        
        # Add bullet points (max 4)
        for i, rec in enumerate(recommendations[:4], 1):
            p = text_frame.add_paragraph()
            p.text = rec
            p.level = 0
            p.font.size = self.body_font_size
            p.space_after = Pt(12)
    
    def _add_swot_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Slide 8: SWOT Summary - Top item from each category."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "SWOT Analysis"
        title.text_frame.paragraphs[0].font.size = self.heading_font_size
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Get SWOT summary
        strategy = data.get("strategy_synthesis", {})
        swot = strategy.get("swot_summary", {})
        
        # Add SWOT items
        swot_items = [
            ("Strength", swot.get("top_strength", "Not identified")),
            ("Weakness", swot.get("top_weakness", "Not identified")),
            ("Opportunity", swot.get("top_opportunity", "Not identified")),
            ("Threat", swot.get("top_threat", "Not identified"))
        ]
        
        for category, item in swot_items:
            p = text_frame.add_paragraph()
            p.text = f"{category}: {item}"
            p.level = 0
            p.font.size = self.body_font_size
            p.space_after = Pt(12)
    
    def _add_next_steps_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Slide 9: Next Steps - Implementation roadmap."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Next Steps"
        title.text_frame.paragraphs[0].font.size = self.heading_font_size
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Generic next steps
        next_steps = [
            "Immediate (0-3 months): Execute quick wins and foundational initiatives",
            "Short-term (3-6 months): Build momentum with strategic priorities",
            "Medium-term (6-12 months): Scale successful initiatives and optimize",
            "Ongoing: Monitor KPIs and adjust strategy based on market feedback"
        ]
        
        for step in next_steps:
            p = text_frame.add_paragraph()
            p.text = step
            p.level = 0
            p.font.size = self.body_font_size
            p.space_after = Pt(12)
    
    def _add_qa_slide(self, prs: Presentation):
        """Slide 10: Q&A placeholder."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Q&A text
        qa_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(8), Inches(1.5)
        )
        qa_frame = qa_box.text_frame
        qa_frame.text = "Questions?"
        qa_para = qa_frame.paragraphs[0]
        qa_para.font.size = Pt(60)
        qa_para.font.bold = True
        qa_para.font.color.rgb = self.primary_color
        qa_para.alignment = PP_ALIGN.CENTER

    # ==================== Comparison Deck ====================

    async def generate_comparison_deck(
        self,
        title: str,
        companies: List[str],
        comparison_data: Dict[str, Any],
        output_path: str,
    ) -> str:
        """Generate a comparison pitch deck with side-by-side slides."""
        logger.info("generating_comparison_deck", companies=companies)

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        tf = tb.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        p.alignment = PP_ALIGN.CENTER

        sub = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
        sf = sub.text_frame
        sf.text = "Comparative Strategic Analysis"
        sp = sf.paragraphs[0]
        sp.font.size = Pt(22)
        sp.font.color.rgb = self.secondary_color
        sp.alignment = PP_ALIGN.CENTER

        date_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        df = date_box.text_frame
        df.text = datetime.now().strftime("%B %Y")
        dp = df.paragraphs[0]
        dp.font.size = Pt(14)
        dp.font.color.rgb = RGBColor(128, 128, 128)
        dp.alignment = PP_ALIGN.CENTER

        # One slide per category
        categories = comparison_data.get("categories", [])
        for cat in categories:
            cat_name = cat.get("name", "")
            rows = cat.get("rows", [])
            if not rows:
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Category title
            ttl = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
            ttf = ttl.text_frame
            ttf.text = cat_name
            tp = ttf.paragraphs[0]
            tp.font.size = Pt(28)
            tp.font.bold = True
            tp.font.color.rgb = self.primary_color

            # Build table
            n_rows = len(rows) + 1  # +1 header
            n_cols = 1 + len(companies)  # metric + companies

            tbl_shape = slide.shapes.add_table(
                n_rows, n_cols,
                Inches(0.5), Inches(1.2), Inches(9), Inches(5.5),
            )
            tbl = tbl_shape.table

            # Header row
            tbl.cell(0, 0).text = "Metric"
            for ci, c in enumerate(companies):
                tbl.cell(0, ci + 1).text = c

            # Style header
            for ci in range(n_cols):
                cell = tbl.cell(0, ci)
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(14)
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 64, 175)

            # Data rows
            for ri, row in enumerate(rows):
                tbl.cell(ri + 1, 0).text = row.get("metric", "")
                for ci in range(len(companies)):
                    tbl.cell(ri + 1, ci + 1).text = str(row.get(f"company_{ci}", "N/A"))

                # Style data cells
                for ci in range(n_cols):
                    cell = tbl.cell(ri + 1, ci)
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(12)

                    # Alternating row colors
                    if ri % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(243, 244, 246)

        # Verdict slide
        verdict = comparison_data.get("verdict", "")
        if verdict:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            vt = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.7))
            vtf = vt.text_frame
            vtf.text = "Verdict"
            vp = vtf.paragraphs[0]
            vp.font.size = Pt(36)
            vp.font.bold = True
            vp.font.color.rgb = self.primary_color
            vp.alignment = PP_ALIGN.CENTER

            vb = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(8), Inches(2))
            vbf = vb.text_frame
            vbf.word_wrap = True
            vbf.text = verdict
            vbp = vbf.paragraphs[0]
            vbp.font.size = Pt(20)
            vbp.alignment = PP_ALIGN.CENTER

        prs.save(output_path)
        logger.info("comparison_deck_generated", path=output_path, slides=len(prs.slides))
        return output_path

